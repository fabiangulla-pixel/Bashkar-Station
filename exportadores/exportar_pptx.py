"""exportadores/exportar_pptx.py — Exportador de presentación PowerPoint.

Genera una presentación .pptx lista para uso académico con:
  - Portada del proyecto
  - Estadísticas del corpus
  - Principales entidades NER por categoría
  - Tópicos detectados
  - Métricas de red
  - Análisis de tono editorial
  - Slide de conclusiones con narrativa Claude
"""

from __future__ import annotations

from pathlib import Path


def exportar_presentacion(
    datos: dict,
    ruta: Path,
    titulo_proyecto: str = "Análisis corpus Estampa",
    investigador: str = "",
    institucion: str = "Instituto Caro y Cuervo",
) -> Path:
    """
    Genera presentación .pptx a partir de los datos del proyecto.

    datos debe contener (todos opcionales):
      - articulos: dict
      - indice_ner_global: dict
      - topicos: dict
      - metricas_red: dict
      - estadisticas_tono: dict
      - narrativa: str
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        raise ImportError("Instala python-pptx: pip install python-pptx")

    ruta = Path(ruta)
    ruta.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Paleta de colores
    C_MORADO  = RGBColor(0x7C, 0x3A, 0xED)
    C_AZUL    = RGBColor(0x38, 0xBD, 0xF8)
    C_VERDE   = RGBColor(0x22, 0xC5, 0x5E)
    C_FONDO   = RGBColor(0x0F, 0x0F, 0x23)
    C_TEXTO   = RGBColor(0xE0, 0xE0, 0xE0)
    C_GRIS    = RGBColor(0x94, 0xA3, 0xB8)

    def _slide_en_blanco(layout_idx: int = 6):
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = C_FONDO
        return slide

    def _caja(slide, x, y, w, h, texto: str, size: int = 18,
               bold: bool = False, color: RGBColor = None, align=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or C_TEXTO

    # ── Slide 1: Portada ─────────────────────────────────────────────────────
    s = _slide_en_blanco()
    _caja(s, 0.5, 1.5, 12, 1.5, titulo_proyecto, size=36, bold=True, color=C_MORADO, align=PP_ALIGN.CENTER)
    _caja(s, 0.5, 3.2, 12, 0.7, "Revista Estampa, Colombia 1930-1940", size=20, color=C_AZUL, align=PP_ALIGN.CENTER)
    if investigador:
        _caja(s, 0.5, 4.2, 12, 0.5, investigador, size=16, color=C_GRIS, align=PP_ALIGN.CENTER)
    if institucion:
        _caja(s, 0.5, 4.9, 12, 0.5, institucion, size=14, color=C_GRIS, align=PP_ALIGN.CENTER)
    _caja(s, 0.5, 6.5, 12, 0.5, "Generado con Bashkar Station", size=11, color=C_GRIS, align=PP_ALIGN.CENTER)

    # ── Slide 2: Estadísticas del corpus ─────────────────────────────────────
    articulos = datos.get("articulos", {})
    n_arts = len(articulos)
    n_pags = sum(len(a.get("paginas", [])) for a in articulos.values())
    n_ents = sum(
        len(ents)
        for ents in datos.get("indice_ner_global", {}).values()
        if isinstance(ents, dict)
    )

    s = _slide_en_blanco()
    _caja(s, 0.5, 0.3, 12, 0.7, "Estadísticas del corpus", size=28, bold=True, color=C_AZUL)
    stats = [
        (str(n_arts), "Artículos"),
        (str(n_pags) if n_pags else "—", "Páginas"),
        (str(n_ents), "Entidades NER"),
        (str(len(datos.get("topicos", {}).get("topicos", []))), "Tópicos"),
    ]
    for i, (val, lab) in enumerate(stats):
        x = 1.0 + i * 3.0
        _caja(s, x, 1.8, 2.5, 1.2, val, size=48, bold=True, color=C_MORADO, align=PP_ALIGN.CENTER)
        _caja(s, x, 3.2, 2.5, 0.6, lab, size=16, color=C_GRIS, align=PP_ALIGN.CENTER)

    # ── Slide 3: Entidades NER ────────────────────────────────────────────────
    indice_ner = datos.get("indice_ner_global", {})
    if indice_ner:
        s = _slide_en_blanco()
        _caja(s, 0.5, 0.3, 12, 0.7, "Principales entidades (NER)", size=28, bold=True, color=C_AZUL)
        cats = list(indice_ner.items())
        col_w = 4.0
        for ci, (cat, ents) in enumerate(cats[:3]):
            col_x = 0.5 + ci * col_w
            _caja(s, col_x, 1.2, col_w - 0.2, 0.5, cat.upper(), size=13, bold=True, color=C_VERDE)
            top = sorted(ents.items(), key=lambda x: -len(x[1]) if isinstance(x[1], (list,set)) else -1)[:8]
            texto_ents = "\n".join(f"• {e[0]}" for e in top)
            _caja(s, col_x, 1.9, col_w - 0.2, 5.0, texto_ents, size=12, color=C_TEXTO)

    # ── Slide 4: Tópicos ──────────────────────────────────────────────────────
    topicos_res = datos.get("topicos", {})
    if topicos_res:
        s = _slide_en_blanco()
        _caja(s, 0.5, 0.3, 12, 0.7, "Tópicos detectados", size=28, bold=True, color=C_AZUL)
        topicos = topicos_res.get("topicos", [])
        # core/topic_engine.py devuelve {id: {palabras, n_docs, nombre}} —un
        # dict, no una lista— y llama a la clave "palabras", no "palabras_top".
        # Iterando la lista sin más, un resultado real del motor recorría las
        # claves (strings) y la diapositiva salía vacía o reventaba.
        if isinstance(topicos, dict):
            topicos = [{"id": k, **v} if isinstance(v, dict) else {"id": k}
                       for k, v in topicos.items()]
        etiquetas = topicos_res.get("etiquetas_llm", {})
        y = 1.5
        for i, top in enumerate(topicos[:8]):
            if not isinstance(top, dict):
                continue
            t_id = top.get("id", i)
            etq = (etiquetas.get(str(t_id)) or etiquetas.get(t_id)
                   or top.get("nombre") or f"Tópico {t_id}")
            palabras = ", ".join(
                (top.get("palabras_top") or top.get("palabras") or [])[:6])
            _caja(s, 0.5, y, 4.0, 0.45, f"● {etq}", size=14, bold=True, color=C_VERDE)
            _caja(s, 4.7, y, 8.0, 0.45, palabras, size=11, color=C_GRIS)
            y += 0.6
            if y > 6.5:
                break

    # ── Slide 5: Red de co-ocurrencias ────────────────────────────────────────
    metricas = datos.get("metricas_red", {})
    if metricas:
        s = _slide_en_blanco()
        _caja(s, 0.5, 0.3, 12, 0.7, "Red de co-ocurrencias", size=28, bold=True, color=C_AZUL)
        items_m = [
            ("Nodos", str(metricas.get("nodos", 0))),
            ("Aristas", str(metricas.get("aristas", 0))),
            ("Densidad", f"{metricas.get('densidad', 0):.3f}"),
            ("Comunidades", str(metricas.get("n_comunidades", 0))),
        ]
        for i, (lab, val) in enumerate(items_m):
            x = 1.0 + i * 3.0
            _caja(s, x, 2.0, 2.5, 1.0, val, size=40, bold=True, color=C_MORADO, align=PP_ALIGN.CENTER)
            _caja(s, x, 3.2, 2.5, 0.5, lab, size=14, color=C_GRIS, align=PP_ALIGN.CENTER)

        top_c = metricas.get("centralidad_top", [])
        if top_c:
            texto_c = "Top entidades por centralidad:\n" + ", ".join(
                f"{n} ({v:.3f})" for n, v in top_c[:5]
            )
            _caja(s, 0.5, 4.5, 12, 1.5, texto_c, size=13, color=C_TEXTO)

    # ── Slide 6: Tono editorial ───────────────────────────────────────────────
    est_tono = datos.get("estadisticas_tono", {})
    if est_tono:
        s = _slide_en_blanco()
        _caja(s, 0.5, 0.3, 12, 0.7, "Tono editorial", size=28, bold=True, color=C_AZUL)
        dist = est_tono.get("distribucion", {})
        y = 1.5
        for tono, pct in sorted(dist.items(), key=lambda x: -x[1]):
            barra = "█" * int(pct * 30)
            _caja(s, 0.5, y, 3.0, 0.45, tono, size=13, bold=True, color=C_VERDE)
            _caja(s, 3.6, y, 7.0, 0.45, f"{barra} {pct:.1%}", size=11, color=C_TEXTO)
            y += 0.55

    # ── Slide 7: Narrativa / Conclusiones ─────────────────────────────────────
    narrativa = datos.get("narrativa", "")
    if narrativa:
        s = _slide_en_blanco()
        _caja(s, 0.5, 0.3, 12, 0.7, "Conclusiones", size=28, bold=True, color=C_AZUL)
        _caja(s, 0.5, 1.3, 12.3, 5.5, narrativa[:900], size=13, color=C_TEXTO)

    prs.save(str(ruta))
    return ruta
